import csv
import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


class PresentationUpdater:
    def __init__(self, csv_file, pptx_file):
        self.csv_file = csv_file
        self.pptx_file = pptx_file
        self.data_dict = {}
        self.presentation = None

    def load_csv_data(self):
        """
        Load data from CSV file into a dictionary.
        """
        with open(self.csv_file, mode="r") as file:
            reader = csv.DictReader(file)
            for row in reader:
                self.data_dict[row["Key"]] = row["Data"]

    def verify_csv_data(self):
        """
        Verify that the CSV data has the correct keys and data types for each field.
        """
        # Define expected data types for each key
        expected_data_types = {
            "img1": "image",
            "title1": "text",
            "desc1": "text",
            "img2": "image",
            "img3": "image",
            "desc2": "text",
        }

        errors = []

        for key, expected_type in expected_data_types.items():
            if key not in self.data_dict:
                errors.append(f"Missing key: {key}")
                continue

            value = self.data_dict[key]

            if expected_type == "image":
                if not os.path.isfile(value):
                    errors.append(
                        f"Expected an image file for key '{key}', but file '{value}' does not exist."
                    )
            elif expected_type == "text":
                if not isinstance(value, str):
                    errors.append(
                        f"Expected a text string for key '{key}', but got value '{value}' of type {type(value)}."
                    )

        return len(errors) == 0, errors

    def verify_pptx_keys(self):
        """
        Verify that the PowerPoint presentation contains all keys from the CSV data.
        """
        pptx_keys = set()

        for slide in self.presentation.slides:
            for shape in slide.shapes:
                if shape.alt_text:
                    pptx_keys.add(shape.alt_text)

        missing_keys = [key for key in self.data_dict if key not in pptx_keys]
        return len(missing_keys) == 0, missing_keys

    def update_text_frame(self, shape):
        """
        Update the text in the shape's text frame if its alt text is in data_dict.
        """
        alt_text = shape.alt_text
        if alt_text in self.data_dict:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.text = self.data_dict[alt_text]

    def update_image(self, shape, slide):
        """
        Update the image in the shape if its alt text is in data_dict.
        """
        alt_text = shape.alt_text
        if alt_text in self.data_dict:
            image_path = self.data_dict[alt_text]
            if os.path.isfile(image_path):
                sp = shape._element
                sp.getparent().remove(sp)
                slide.shapes.add_picture(
                    image_path, shape.left, shape.top, shape.width, shape.height
                )

    def update_presentation(self):
        """
        Update the PowerPoint presentation based on the CSV data.
        """
        self.load_csv_data()

        # Verify CSV data
        valid, errors = self.verify_csv_data()
        if not valid:
            print("CSV data verification failed with the following errors:")
            for error in errors:
                print(f" - {error}")
            return

        self.presentation = Presentation(self.pptx_file)

        # Verify PPTX keys
        pptx_valid, missing_keys = self.verify_pptx_keys()
        if not pptx_valid:
            print("PPTX key verification failed. Missing keys:")
            for key in missing_keys:
                print(f" - {key}")
            return

        # Update presentation
        for slide in self.presentation.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    self.update_text_frame(shape)
                elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    self.update_image(shape, slide)

        updated_pptx_file = "updated_presentation.pptx"
        self.presentation.save(updated_pptx_file)
        print(f"Presentation updated successfully. Saved as {updated_pptx_file}.")


def main():
    """
    Main function to update the PowerPoint presentation.
    """
    csv_file = "data.csv"
    pptx_file = "my_ppt.pptx"

    updater = PresentationUpdater(csv_file, pptx_file)
    updater.update_presentation()


if __name__ == "__main__":
    main()
